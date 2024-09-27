from pptx import Presentation
from pptx.util import Inches
import requests
import os
from io import BytesIO
from PIL import Image
import logging 

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Function to create a PPTX presentation from given data
def create_slide(data, output_path="presentation.pptx"):
    # Create a PowerPoint presentation object
    prs = Presentation()

    # Iterate through the data to create slides
    for item in data:
        # Add a slide with a title and content layout
        slide_layout = prs.slide_layouts[1]  # Using layout with title and content
        slide = prs.slides.add_slide(slide_layout)

        # Add the title (text field)
        title = slide.shapes.title
        title.text = item['text']

        # Add the voice (as content in the slide)
        content_box = slide.shapes.placeholders[1]
        content_box.text = item['voice']


        # Get slide width (the total width of the slide)
        slide_width = prs.slide_width

        # Adjust content box to fit below the heading
        content_box.top = int(title.top + title.height)
        content_box.width = int(slide_width / 2)  # Set width to half of the slide width
        content_box.left = 0

         # Add the image from the URL and place it at the bottom right of the slide
        if 'image_url' in item and item['image_url']:
            image_url = item['image_url']
            try:
                # Fetch the image from the URL
                response = requests.get(image_url)
                img = Image.open(BytesIO(response.content))

                # Save the image to a temporary buffer and add to the slide
                img_buffer = BytesIO()
                img.save(img_buffer, format="PNG")
                img_buffer.seek(0)

                # Calculate the position for the image (center-right)
                img_left = int(slide_width / 2)  # Start from the middle of the slide
                img_top = int(prs.slide_height / 2 - Inches(1.5))  # Vertically centered

                # Add the image to the slide, positioned at the center-right
                slide.shapes.add_picture(img_buffer, Inches(6), img_top, width=Inches(3))
            except Exception as e:
                logger.info(f"Error adding image from {image_url}: {e}")

    # Save the PowerPoint presentation
    prs.save(output_path)
    return os.path.abspath(output_path)